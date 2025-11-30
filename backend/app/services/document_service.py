from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
import io
from app.utils.logger import get_logger
from typing import List, Dict

logger = get_logger(__name__)

class DocumentService:
    @staticmethod
    def create_word_document(project_data: dict) -> bytes:
        """
        Create professionally formatted Word document
        Uses latest refined content from each section
        """
        try:
            doc = Document()
            
            # ===== Title Page =====
            title = doc.add_heading(project_data['title'], 0)
            title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            
            # Subtitle/Topic
            if project_data.get('topic'):
                topic_para = doc.add_paragraph(project_data['topic'])
                topic_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                topic_para.runs[0].font.size = Pt(14)
                topic_para.runs[0].italic = True
            
            # Creation date
            date_para = doc.add_paragraph(f"Generated: {project_data.get('created_at', 'N/A')}")
            date_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            date_para.runs[0].font.size = Pt(10)
            date_para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            doc.add_page_break()
            
            # ===== Table of Contents =====
            toc_heading = doc.add_heading('Table of Contents', 1)
            for idx, section in enumerate(project_data['sections'], 1):
                toc_item = doc.add_paragraph(f"{idx}. {section['title']}")
                toc_item.style = 'List Number'
                toc_item.paragraph_format.left_indent = Pt(20)
            
            doc.add_page_break()
            
            # ===== Content Sections =====
            for section in sorted(project_data['sections'], key=lambda x: x.get('order', 0)):
                # Section heading
                heading = doc.add_heading(section['title'], level=1)
                
                # Section content (use latest refined content)
                content = section.get('content', '')
                
                # Split into paragraphs
                paragraphs = content.split('\n\n')
                
                for para_text in paragraphs:
                    if para_text.strip():
                        # Check if it's a bullet point
                        if para_text.strip().startswith(('•', '-', '*')):
                            para = doc.add_paragraph(para_text.strip().lstrip('•-* '), style='List Bullet')
                        else:
                            para = doc.add_paragraph(para_text.strip())
                        
                        # Format paragraph
                        para.paragraph_format.line_spacing = 1.5
                        para.paragraph_format.space_after = Pt(12)
                        
                        # Format text
                        for run in para.runs:
                            run.font.name = 'Calibri'
                            run.font.size = Pt(11)
                
                # Add spacing after section
                doc.add_paragraph()
            
            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            logger.info(f"Word document created: {project_data['title']}")
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Word document creation error: {str(e)}")
            raise
    
    @staticmethod
    def create_powerpoint(project_data: dict) -> bytes:
        """
        Create professionally formatted PowerPoint presentation
        Uses latest refined content from each slide
        """
        try:
            prs = Presentation()
            prs.slide_width = PptxInches(10)
            prs.slide_height = PptxInches(7.5)
            
            # ===== Title Slide =====
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = project_data['title']
            subtitle_text = project_data.get('topic', '')
            if project_data.get('created_at'):
                subtitle_text += f"\n{project_data['created_at']}"
            subtitle_shape.text = subtitle_text
            
            # ===== Content Slides =====
            for section in sorted(project_data['sections'], key=lambda x: x.get('order', 0)):
                # Use bullet slide layout
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Slide title
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                
                # Slide content
                content_box = slide.placeholders[1]
                text_frame = content_box.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                
                # Get latest content
                content = section.get('content', '')
                
                # Parse content into bullet points
                lines = content.split('\n')
                bullet_points = []
                
                for line in lines:
                    line = line.strip()
                    if line:
                        # Remove existing bullet markers
                        line = line.lstrip('•-* ')
                        if line:
                            bullet_points.append(line)
                
                # Add bullet points (limit to 6 per slide for readability)
                for point in bullet_points[:6]:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = PptxPt(18)
                    p.space_after = PptxPt(12)
            
            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            logger.info(f"PowerPoint created: {project_data['title']}")
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"PowerPoint creation error: {str(e)}")
            raise
